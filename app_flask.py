"""Briefly - Flask backend (mounted under /api/*).

Improvements over original:
- Document Maker endpoint that produces real .docx, .xlsx, .pdf, .pptx, .txt, .md files
  using Python libraries (no LLM round-trip required for the file build itself).
- Auto-create conversation when user sends first message (frontend already supports;
  backend just makes sure conversation_id is optional on /api/chat).
- New /api/config endpoint to expose public reCAPTCHA key to the frontend.
"""
import os
import io
import json
import re
import datetime as dt
from flask import Flask, request, jsonify, session, send_file, Response, stream_with_context, render_template
from werkzeug.utils import secure_filename
from dotenv import load_dotenv
import requests
import PyPDF2
import docx
from google import genai

from backend.src.database import (
    init_db, register_user, login_user, create_conversation,
    get_user_conversations, get_conversation_messages, save_message,
    delete_conversation, get_user_profile, update_user_profile, change_password,
    get_db,
    save_analyzed_document, list_analyzed_documents, get_analyzed_document, delete_analyzed_document,
    list_drafts, get_draft, create_draft, update_draft, delete_draft,
)
from backend.src.rag_system import RAGSystem

load_dotenv(os.path.join(os.path.dirname(__file__), '.env'))

app = Flask(__name__)
@app.route("/")
def home():
    return render_template("index.html")

app.secret_key = os.getenv('SECRET_KEY', 'change-me')
app.config['UPLOAD_FOLDER'] = os.path.join(os.path.dirname(__file__), 'uploads')
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'doc', 'docx'}

RECAPTCHA_SITE_KEY = os.getenv('RECAPTCHA_SITE_KEY', '')
RECAPTCHA_SECRET_KEY = os.getenv('RECAPTCHA_SECRET_KEY', '')

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
init_db()

GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
GEMINI_MODEL = os.getenv('GEMINI_MODEL', 'gemini-2.5-flash')
if not GEMINI_API_KEY:
    raise ValueError('GEMINI_API_KEY missing in .env')

client = genai.Client(api_key=GEMINI_API_KEY)
rag_system = RAGSystem(GEMINI_API_KEY, GEMINI_MODEL)


def _load_simple_responses():
    try:
        path = os.path.join(os.path.dirname(__file__), 'backend/src', 'simple_responses.json')
        with open(path) as f:
            return json.load(f)['simple_responses']
    except Exception:
        return {}

simple_responses = _load_simple_responses()


def verify_recaptcha(token):
    # Dev bypass — when DEBUG=True we don't block on recaptcha so testing tools work
    if (os.getenv('DEBUG', 'False').lower() == 'true') and (token in (None, '', 'dummy', 'dev-bypass')):
        return True
    if not token:
        return False
    if not RECAPTCHA_SECRET_KEY:
        return True  # If not configured, don't block
    try:
        r = requests.post(
            'https://www.google.com/recaptcha/api/siteverify',
            data={'secret': RECAPTCHA_SECRET_KEY, 'response': token}, timeout=5,
        )
        return r.json().get('success', False)
    except Exception:
        return False


def allowed_file(fn):
    return '.' in fn and fn.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def preprocess_text(t):
    return ' '.join((t or '').split()).replace('\x00', '')


def read_file_content(filepath):
    try:
        if filepath.endswith('.txt'):
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return preprocess_text(f.read())
        if filepath.endswith('.pdf'):
            text = ''
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for p in reader.pages:
                    text += (p.extract_text() or '') + '\n'
            return preprocess_text(text)
        if filepath.endswith('.docx'):
            d = docx.Document(filepath)
            text = '\n'.join(p.text for p in d.paragraphs)
            for table in d.tables:
                for row in table.rows:
                    text += '\n' + ' | '.join(c.text for c in row.cells)
            return preprocess_text(text)
        return 'Unsupported file format'
    except Exception as e:
        return f'Error reading file: {e}'


# =========== AUTH ===========
@app.route('/api/config')
def api_config():
    return jsonify({'recaptcha_site_key': RECAPTCHA_SITE_KEY})


@app.route('/api/register', methods=['POST'])
def register():
    data = request.json or {}
    name = (data.get('name') or '').strip()
    email = (data.get('email') or '').strip()
    password = (data.get('password') or '').strip()
    if not email or not password:
        return jsonify({'success': False, 'error': 'Email and password required'}), 400
    if len(password) < 6:
        return jsonify({'success': False, 'error': 'Password must be at least 6 characters'}), 400
    if not verify_recaptcha(data.get('recaptcha_response')):
        return jsonify({'success': False, 'error': 'reCAPTCHA verification failed'}), 400
    if register_user(email, password, name, email):
        return jsonify({'success': True})
    return jsonify({'success': False, 'error': 'Email already exists'}), 400


@app.route('/api/login', methods=['POST'])
def login():
    data = request.json or {}
    username = (data.get('username') or '').strip()
    password = (data.get('password') or '').strip()
    if not username or not password:
        return jsonify({'success': False, 'error': 'Username and password required'}), 400
    if not verify_recaptcha(data.get('recaptcha_response')):
        return jsonify({'success': False, 'error': 'reCAPTCHA verification failed'}), 400
    user_id = login_user(username, password)
    if user_id:
        session['user_id'] = user_id
        session['username'] = username
        return jsonify({'success': True})
    return jsonify({'success': False, 'error': 'Invalid credentials'}), 401


@app.route('/api/logout', methods=['POST'])
def logout():
    session.clear()
    return jsonify({'success': True})


@app.route('/api/user', methods=['GET'])
def get_user():
    if 'user_id' in session:
        return jsonify({'success': True, 'user_id': session['user_id'],
                        'username': session['username'],
                        'profile': get_user_profile(session['user_id'])})
    return jsonify({'success': False, 'error': 'Not logged in'}), 401


# =========== CONVERSATIONS ===========
@app.route('/api/conversations', methods=['GET'])
def get_conversations():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    return jsonify({'success': True, 'conversations': get_user_conversations(session['user_id'])})


@app.route('/api/conversations/new', methods=['POST'])
def new_conversation():
    data = request.json or {}
    title = data.get('title') or 'New Chat'
    user_id = session.get('user_id', 0)
    cid = create_conversation(user_id, title)
    return jsonify({'success': True, 'conversation_id': cid})


@app.route('/api/conversations/<int:cid>/messages', methods=['GET'])
def get_messages(cid):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    conn = get_db(); cur = conn.cursor()
    cur.execute('SELECT user_id FROM conversations WHERE id = ?', (cid,))
    row = cur.fetchone(); conn.close()
    if not row or row[0] != session['user_id']:
        return jsonify({'success': False, 'error': 'Access denied'}), 403
    return jsonify({'success': True, 'messages': get_conversation_messages(cid)})


@app.route('/api/conversations/<int:cid>/delete', methods=['POST'])
def delete_conv(cid):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    conn = get_db(); cur = conn.cursor()
    cur.execute('SELECT user_id FROM conversations WHERE id = ?', (cid,))
    row = cur.fetchone(); conn.close()
    if not row or row[0] != session['user_id']:
        return jsonify({'success': False, 'error': 'Access denied'}), 403
    delete_conversation(cid)
    return jsonify({'success': True})


# =========== CHAT ===========
@app.route('/api/chat', methods=['POST'])
def chat():
    data = request.json or {}
    user_message = (data.get('message') or '').strip()
    conversation_id = data.get('conversation_id')
    if not user_message:
        return jsonify({'success': False, 'error': 'Message cannot be empty'}), 400

    # Auto-create conversation if not provided (only for logged-in users)
    if not conversation_id:
        user_id = session.get('user_id', 0)
        title = user_message[:50] + ('...' if len(user_message) > 50 else '')
        conversation_id = create_conversation(user_id, title)

    if 'user_id' in session:
        conn = get_db(); cur = conn.cursor()
        cur.execute('SELECT user_id FROM conversations WHERE id = ?', (conversation_id,))
        row = cur.fetchone(); conn.close()
        if row and row[0] != session['user_id'] and row[0] != 0:
            return jsonify({'success': False, 'error': 'Access denied'}), 403

    lower = user_message.lower().strip()
    if lower in simple_responses:
        reply = simple_responses[lower]
    else:
        previous = get_conversation_messages(conversation_id)
        reply = rag_system.generate_response_with_rag(user_message, previous)

    save_message(conversation_id, 'user', user_message)
    save_message(conversation_id, 'assistant', reply)
    return jsonify({'success': True, 'message': reply, 'conversation_id': conversation_id})


# =========== PROFILE ===========
@app.route('/api/profile/update', methods=['POST'])
def update_profile():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    d = request.json or {}
    update_user_profile(session['user_id'], d.get('name'), d.get('mobile'), d.get('dob'))
    return jsonify({'success': True})


@app.route('/api/profile/change-password', methods=['POST'])
def change_password_endpoint():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    d = request.json or {}
    cur_pwd = d.get('current_password') or ''
    new_pwd = d.get('new_password') or ''
    if not cur_pwd or not new_pwd:
        return jsonify({'success': False, 'error': 'Current and new password required'}), 400
    if not login_user(session['username'], cur_pwd):
        return jsonify({'success': False, 'error': 'Current password is incorrect'}), 400
    if len(new_pwd) < 6:
        return jsonify({'success': False, 'error': 'New password must be at least 6 characters'}), 400
    change_password(session['user_id'], new_pwd)
    return jsonify({'success': True})


# =========== DOCUMENT ANALYZER ===========
@app.route('/api/upload-document', methods=['POST'])
def upload_document():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    if 'file' not in request.files:
        return jsonify({'success': False, 'error': 'No file provided'}), 400
    f = request.files['file']
    if not f.filename or not allowed_file(f.filename):
        return jsonify({'success': False, 'error': 'Invalid file type. Allowed: txt, pdf, doc, docx'}), 400
    fn = secure_filename(f.filename)
    fp = os.path.join(app.config['UPLOAD_FOLDER'], fn)
    f.save(fp)
    try:
        content = read_file_content(fp)
    finally:
        try: os.remove(fp)
        except Exception: pass

    summary_resp = client.models.generate_content(
        model=GEMINI_MODEL,
        contents=f"Summarise this document concisely in 4-6 sentences:\n\n{content[:8000]}"
    )
    summary = (summary_resp.text or '').strip()
    topics_resp = client.models.generate_content(
        model=GEMINI_MODEL,
        contents=f"Extract 3-5 short topic labels (comma-separated, no numbering) from this summary:\n\n{summary}"
    )
    topics = (topics_resp.text or '').strip()

    # Persist for history
    display_name = (request.form.get('name') or '').strip() or fn
    doc_id = save_analyzed_document(session['user_id'], fn, display_name, summary, topics, content)

    return jsonify({'success': True, 'id': doc_id, 'summary': summary, 'topics': topics,
                    'filename': fn, 'display_name': display_name, 'content': content})


@app.route('/api/document-chat', methods=['POST'])
def document_chat():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    d = request.json or {}
    msg = (d.get('message') or '').strip()
    if not msg:
        return jsonify({'success': False, 'error': 'Message cannot be empty'}), 400
    prompt = f"""Based on this document, answer the user's question.

DOCUMENT SUMMARY:
{d.get('document_summary','')}

DOCUMENT CONTENT:
{(d.get('document_content','') or '')[:12000]}

USER QUESTION: {msg}

Answer based on the document. If the term appears in the document, you may add a brief general-knowledge explanation. If it's clearly off-topic, say so politely."""
    resp = client.models.generate_content(model=GEMINI_MODEL, contents=prompt)
    return jsonify({'success': True, 'response': (resp.text or '').strip()})


# =========== DOCUMENT MAKER (NEW) ===========
DOC_TEMPLATES = {
    'bill':    {'label': 'Bill / Invoice',          'icon': 'fa-file-invoice-dollar'},
    'resume':  {'label': 'Resume / CV',             'icon': 'fa-id-card'},
    'project_architecture': {'label': 'Project Architecture', 'icon': 'fa-diagram-project'},
    'school_project':  {'label': 'School Project',  'icon': 'fa-school'},
    'college_project': {'label': 'College Project', 'icon': 'fa-graduation-cap'},
    'thesis':  {'label': 'Thesis / Dissertation',   'icon': 'fa-book'},
    'report':  {'label': 'Report',                  'icon': 'fa-file-lines'},
    'letter':  {'label': 'Letter',                  'icon': 'fa-envelope-open-text'},
    'meeting_minutes': {'label': 'Meeting Minutes', 'icon': 'fa-users'},
    'business_plan':  {'label': 'Business Plan',    'icon': 'fa-briefcase'},
    'custom':  {'label': 'Custom Document',         'icon': 'fa-file-pen'},
}

DOC_FORMATS = ['docx', 'pdf', 'xlsx', 'pptx', 'txt', 'md']


@app.route('/api/document-maker/types', methods=['GET'])
def doc_maker_types():
    return jsonify({'success': True, 'types': DOC_TEMPLATES, 'formats': DOC_FORMATS})


def _safe_name(s):
    s = re.sub(r'[^A-Za-z0-9_\- ]+', '', s or '').strip().replace(' ', '_') or 'document'
    return s[:60]


def _expand_with_llm(doc_type, data_dict):
    """OPTIONAL: lightly enrich free-text fields. Most of the heavy lifting
    is done by Python libs — LLM is only used to expand a single 'content'
    free-text field into structured sections when needed."""
    raw = data_dict.get('content', '') or ''
    if not raw or len(raw) > 200:
        return raw  # already detailed enough, skip LLM
    try:
        prompt = (
            f"Expand the following short prompt into a well-structured {DOC_TEMPLATES.get(doc_type,{}).get('label','document')} "
            f"of about 300-500 words. Return plain text with clear paragraph breaks; do NOT add markdown.\n\nPrompt: {raw}"
        )
        resp = client.models.generate_content(model=GEMINI_MODEL, contents=prompt)
        return (resp.text or raw).strip()
    except Exception:
        return raw


# ----- Builders -----

# ====== Resume templates (multi-template, simple black text) ======
RESUME_TEMPLATES = ['classic', 'modern', 'compact']


def _docx_para_with_hyperlinks(d, text, *, bold=False, size=None):
    """Add a paragraph that auto-detects emails/phones/URLs and renders them as
    hyperlinks (blue + underline). Everything else is plain black text."""
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    p = d.add_paragraph()
    # Combined detector
    pattern = re.compile(r'(https?://[^\s]+|[\w._%+-]+@[\w.-]+\.[A-Za-z]{2,}|\+?\d[\d\s().-]{6,}\d)')
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            run = p.add_run(text[pos:m.start()])
            run.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)
            if size: run.font.size = Pt(size)
        target = m.group(0)
        if '@' in target and ':' not in target:
            href = 'mailto:' + target
        elif target.startswith('+') or (target[0].isdigit() and not target.startswith('http')):
            href = 'tel:' + re.sub(r'[\s().-]', '', target)
        else:
            href = target
        _add_hyperlink(p, href, target, bold=bold, size=size)
        pos = m.end()
    if pos < len(text):
        run = p.add_run(text[pos:])
        run.bold = bold
        run.font.color.rgb = RGBColor(0, 0, 0)
        if size: run.font.size = Pt(size)
    return p


def _add_hyperlink(paragraph, url, text, bold=False, size=None):
    """Insert a hyperlink run into a python-docx paragraph (blue + underline)."""
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    part = paragraph.part
    r_id = part.relate_to(url, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
    hyperlink = OxmlElement('w:hyperlink')
    hyperlink.set(qn('r:id'), r_id)
    new_run = OxmlElement('w:r')
    rPr = OxmlElement('w:rPr')
    color = OxmlElement('w:color'); color.set(qn('w:val'), '0563C1'); rPr.append(color)
    u = OxmlElement('w:u'); u.set(qn('w:val'), 'single'); rPr.append(u)
    if bold:
        b = OxmlElement('w:b'); rPr.append(b)
    if size:
        sz = OxmlElement('w:sz'); sz.set(qn('w:val'), str(size * 2)); rPr.append(sz)
    new_run.append(rPr)
    txt = OxmlElement('w:t'); txt.text = text; txt.set(qn('xml:space'), 'preserve')
    new_run.append(txt)
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)


def _build_resume_docx(d, data):
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    template = (data.get('template') or 'classic').lower()
    if template not in RESUME_TEMPLATES:
        template = 'classic'

    name = data.get('full_name', '').strip()
    title = data.get('title', '').strip() if data.get('title') and data.get('title') != 'Resume' else (data.get('headline') or '').strip()
    contact_bits = [b for b in [data.get('email'), data.get('phone'), data.get('location'), data.get('website')] if b]

    # ---- Header ----
    if template == 'modern':
        # Centered header
        p = d.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if name:
            r = p.add_run(name); r.bold = True; r.font.size = Pt(20); r.font.color.rgb = RGBColor(0, 0, 0)
        if title:
            p2 = d.add_paragraph(); p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r = p2.add_run(title); r.font.size = Pt(11); r.font.color.rgb = RGBColor(0, 0, 0)
        if contact_bits:
            p3 = d.add_paragraph(); p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
            joined = '  |  '.join(contact_bits)
            _docx_inline_contacts(p3, joined, size=10)
    else:
        # Classic / compact — left-aligned
        if name:
            p = d.add_paragraph()
            r = p.add_run(name); r.bold = True
            r.font.size = Pt(18 if template == 'classic' else 16)
            r.font.color.rgb = RGBColor(0, 0, 0)
        if title:
            p = d.add_paragraph()
            r = p.add_run(title); r.font.size = Pt(11); r.font.color.rgb = RGBColor(0, 0, 0)
        if contact_bits:
            p = d.add_paragraph()
            _docx_inline_contacts(p, ' • '.join(contact_bits), size=10)

    section_size = 12 if template == 'compact' else 13

    def section(heading):
        p = d.add_paragraph()
        r = p.add_run(heading.upper() if template == 'compact' else heading)
        r.bold = True; r.font.size = Pt(section_size); r.font.color.rgb = RGBColor(0, 0, 0)

    def body_para(text, bullet=False):
        if bullet:
            p = d.add_paragraph(style='List Bullet')
        else:
            p = d.add_paragraph()
        _docx_inline_contacts(p, text, size=10 if template == 'compact' else 11)

    if data.get('summary'):
        section('Profile')
        body_para(data['summary'])
    if data.get('experience'):
        section('Work Experience')
        for line in str(data['experience']).split('\n'):
            line = line.strip()
            if not line: continue
            body_para(line, bullet=(template != 'classic'))
    if data.get('education'):
        section('Education')
        for line in str(data['education']).split('\n'):
            line = line.strip()
            if line: body_para(line, bullet=(template != 'classic'))
    if data.get('skills'):
        section('Technical Skills')
        body_para(data['skills'])
    if data.get('projects'):
        section('Projects')
        for line in str(data['projects']).split('\n'):
            line = line.strip()
            if line: body_para(line, bullet=(template != 'classic'))
    if data.get('hobbies'):
        section('Hobbies')
        body_para(data['hobbies'])
    if data.get('languages'):
        section('Languages Known')
        body_para(data['languages'])


def _docx_inline_contacts(p, text, size=11):
    """Same as _docx_para_with_hyperlinks but uses an existing paragraph."""
    from docx.shared import Pt, RGBColor
    pattern = re.compile(r'(https?://[^\s]+|[\w._%+-]+@[\w.-]+\.[A-Za-z]{2,}|\+?\d[\d\s().-]{6,}\d)')
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            run = p.add_run(text[pos:m.start()])
            run.font.color.rgb = RGBColor(0, 0, 0)
            if size: run.font.size = Pt(size)
        target = m.group(0)
        if '@' in target and not target.startswith('http'):
            href = 'mailto:' + target
        elif (target.startswith('+') or target[0].isdigit()) and not target.startswith('http'):
            href = 'tel:' + re.sub(r'[\s().-]', '', target)
        else:
            href = target
        _add_hyperlink(p, href, target, size=size)
        pos = m.end()
    if pos < len(text):
        run = p.add_run(text[pos:])
        run.font.color.rgb = RGBColor(0, 0, 0)
        if size: run.font.size = Pt(size)


def _build_resume_pdf(story, data, ss):
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    from reportlab.platypus import Paragraph, Spacer, HRFlowable

    template = (data.get('template') or 'classic').lower()
    if template not in RESUME_TEMPLATES:
        template = 'classic'

    BLACK = colors.HexColor('#000000')
    LINK_BLUE = colors.HexColor('#0563C1')

    name_style = ParagraphStyle('rname', parent=ss['Title'], textColor=BLACK,
                                fontSize=22 if template != 'compact' else 18,
                                alignment=1 if template == 'modern' else 0,
                                spaceAfter=2, leading=24)
    title_style2 = ParagraphStyle('rtitle', parent=ss['BodyText'], textColor=BLACK,
                                 fontSize=11, alignment=1 if template == 'modern' else 0,
                                 spaceAfter=6)
    contact_style = ParagraphStyle('rcontact', parent=ss['BodyText'], textColor=BLACK,
                                   fontSize=10, alignment=1 if template == 'modern' else 0,
                                   spaceAfter=10)
    section_style = ParagraphStyle('rsection', parent=ss['BodyText'], textColor=BLACK,
                                   fontSize=12 if template == 'compact' else 13,
                                   fontName='Helvetica-Bold', spaceBefore=10, spaceAfter=4)
    body_style = ParagraphStyle('rbody', parent=ss['BodyText'], textColor=BLACK,
                                fontSize=10 if template == 'compact' else 11,
                                spaceAfter=4, leading=14)

    def linkify(text):
        # turn emails / phones / urls into <a> tags (blue underline)
        def repl(m):
            t = m.group(0)
            if '@' in t and not t.startswith('http'):
                href = 'mailto:' + t
            elif (t.startswith('+') or t[0].isdigit()) and not t.startswith('http'):
                href = 'tel:' + re.sub(r'[\s().-]', '', t)
            else:
                href = t
            return f'<font color="#0563C1"><u><a href="{href}">{t}</a></u></font>'
        return re.sub(r'https?://[^\s<]+|[\w._%+-]+@[\w.-]+\.[A-Za-z]{2,}|\+?\d[\d\s().-]{6,}\d', repl, text)

    name = data.get('full_name', '').strip()
    headline = (data.get('headline') or '').strip()
    contact_bits = [b for b in [data.get('email'), data.get('phone'), data.get('location'), data.get('website')] if b]

    if name:
        story.append(Paragraph(f"<b>{name}</b>", name_style))
    if headline:
        story.append(Paragraph(headline, title_style2))
    if contact_bits:
        sep = '  |  ' if template == 'modern' else ' &bull; '
        story.append(Paragraph(linkify(sep.join(contact_bits)), contact_style))

    if template == 'modern':
        story.append(HRFlowable(width='100%', thickness=0.5, color=BLACK, spaceBefore=2, spaceAfter=8))

    def add_section(title, value, bulletize=True):
        if not value: return
        story.append(Paragraph(title.upper() if template == 'compact' else title, section_style))
        if isinstance(value, str) and '\n' in value and bulletize:
            for line in value.split('\n'):
                line = line.strip()
                if not line: continue
                if template == 'classic':
                    story.append(Paragraph(linkify(line), body_style))
                else:
                    story.append(Paragraph('• ' + linkify(line), body_style))
        else:
            story.append(Paragraph(linkify(str(value)), body_style))

    add_section('Profile', data.get('summary'), bulletize=False)
    add_section('Work Experience', data.get('experience'))
    add_section('Education', data.get('education'))
    add_section('Technical Skills', data.get('skills'), bulletize=False)
    add_section('Projects', data.get('projects'))
    add_section('Hobbies', data.get('hobbies'), bulletize=False)
    add_section('Languages Known', data.get('languages'), bulletize=False)


# ====== End resume templates ======

def _build_docx(doc_type, data, out):
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    d = Document()
    if doc_type != 'resume':
        title = data.get('title') or DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document')
        h = d.add_heading(title, level=0)
        for run in h.runs:
            run.font.color.rgb = RGBColor(0x1B, 0x5E, 0x20)

    if doc_type == 'resume':
        _build_resume_docx(d, data)
    elif doc_type == 'bill':
        d.add_paragraph(f"Invoice #: {data.get('invoice_number','INV-001')}")
        d.add_paragraph(f"Date: {data.get('date', dt.date.today().isoformat())}")
        d.add_paragraph(f"Bill To: {data.get('client_name','')}")
        if data.get('client_address'): d.add_paragraph(data['client_address'])
        d.add_paragraph(' ')
        items = data.get('items') or []
        tbl = d.add_table(rows=1, cols=4); tbl.style = 'Light Grid Accent 1'
        hdr = tbl.rows[0].cells
        hdr[0].text, hdr[1].text, hdr[2].text, hdr[3].text = 'Item', 'Qty', 'Price', 'Total'
        total = 0.0
        for it in items:
            row = tbl.add_row().cells
            qty = float(it.get('qty', 1)); price = float(it.get('price', 0))
            line = qty * price; total += line
            row[0].text = str(it.get('name', ''))
            row[1].text = str(qty); row[2].text = f"{price:.2f}"; row[3].text = f"{line:.2f}"
        d.add_paragraph(' ')
        p = d.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        run = p.add_run(f"Total: {data.get('currency','$')}{total:.2f}"); run.bold = True
        if data.get('notes'): d.add_paragraph(data['notes'])
    else:
        # Generic doc — use sections if provided
        content = _expand_with_llm(doc_type, data) if data.get('content') else ''
        sections = data.get('sections') or []
        if sections:
            for sec in sections:
                if sec.get('heading'): d.add_heading(sec['heading'], level=1)
                if sec.get('body'): d.add_paragraph(sec['body'])
        if content:
            for para in content.split('\n\n'):
                if para.strip(): d.add_paragraph(para.strip())

    d.save(out)


def _build_pdf(doc_type, data, out):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle

    doc = SimpleDocTemplate(out, pagesize=A4, leftMargin=1.8*cm, rightMargin=1.8*cm,
                            topMargin=2*cm, bottomMargin=2*cm)
    ss = getSampleStyleSheet()
    title_style = ParagraphStyle('t', parent=ss['Title'], textColor=colors.HexColor('#1B5E20'))
    body = ss['BodyText']; body.spaceAfter = 6
    h1 = ParagraphStyle('h1', parent=ss['Heading2'], textColor=colors.HexColor('#1B5E20'))
    story = []
    if doc_type != 'resume':
        title = data.get('title') or DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document')
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 12))

    if doc_type == 'bill':
        story.append(Paragraph(f"Invoice #: {data.get('invoice_number','INV-001')}", body))
        story.append(Paragraph(f"Date: {data.get('date', dt.date.today().isoformat())}", body))
        story.append(Paragraph(f"Bill To: {data.get('client_name','')}", body))
        if data.get('client_address'): story.append(Paragraph(data['client_address'], body))
        story.append(Spacer(1, 10))
        rows = [['Item', 'Qty', 'Price', 'Total']]; total = 0.0
        for it in (data.get('items') or []):
            qty = float(it.get('qty', 1)); price = float(it.get('price', 0))
            line = qty * price; total += line
            rows.append([it.get('name',''), f"{qty:g}", f"{price:.2f}", f"{line:.2f}"])
        rows.append(['', '', 'Total', f"{data.get('currency','$')}{total:.2f}"])
        t = Table(rows, colWidths=[7*cm, 2*cm, 3*cm, 3*cm])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1B5E20')),
            ('TEXTCOLOR', (0,0), (-1,0), colors.white),
            ('GRID', (0,0), (-1,-1), 0.4, colors.grey),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('BACKGROUND', (0,-1), (-1,-1), colors.HexColor('#E8F5E9')),
        ]))
        story.append(t)
        if data.get('notes'):
            story.append(Spacer(1, 12)); story.append(Paragraph(data['notes'], body))
    elif doc_type == 'resume':
        _build_resume_pdf(story, data, ss)
    else:
        content = _expand_with_llm(doc_type, data) if data.get('content') else ''
        for sec in (data.get('sections') or []):
            if sec.get('heading'): story.append(Paragraph(sec['heading'], h1))
            if sec.get('body'): story.append(Paragraph(sec['body'].replace('\n','<br/>'), body))
        if content:
            for para in content.split('\n\n'):
                if para.strip(): story.append(Paragraph(para.strip().replace('\n','<br/>'), body))

    doc.build(story)


def _build_xlsx(doc_type, data, out):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    wb = Workbook(); ws = wb.active
    sheet_title = re.sub(r'[\\/?*\[\]:]', '-', DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document'))
    ws.title = sheet_title[:30]
    title = data.get('title') or DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document')
    ws['A1'] = title
    ws['A1'].font = Font(size=16, bold=True, color='1B5E20')
    ws.merge_cells('A1:D1')
    ws['A1'].alignment = Alignment(horizontal='center')

    if doc_type == 'bill':
        ws['A3'] = 'Invoice #'; ws['B3'] = data.get('invoice_number','INV-001')
        ws['A4'] = 'Date';      ws['B4'] = data.get('date', dt.date.today().isoformat())
        ws['A5'] = 'Bill To';   ws['B5'] = data.get('client_name','')
        headers = ['Item','Qty','Price','Total']
        for i, h in enumerate(headers, 1):
            c = ws.cell(row=7, column=i, value=h)
            c.font = Font(bold=True, color='FFFFFF')
            c.fill = PatternFill('solid', fgColor='1B5E20')
        r = 8; total = 0.0
        for it in (data.get('items') or []):
            qty = float(it.get('qty',1)); price = float(it.get('price',0))
            line = qty*price; total += line
            ws.cell(row=r, column=1, value=it.get('name',''))
            ws.cell(row=r, column=2, value=qty)
            ws.cell(row=r, column=3, value=price)
            ws.cell(row=r, column=4, value=line)
            r += 1
        ws.cell(row=r, column=3, value='Total').font = Font(bold=True)
        ws.cell(row=r, column=4, value=total).font = Font(bold=True)
    else:
        # Generic: write key/value pairs
        r = 3
        for k, v in data.items():
            if k in ('title','sections') or v is None: continue
            ws.cell(row=r, column=1, value=str(k).replace('_',' ').title()).font = Font(bold=True)
            ws.cell(row=r, column=2, value=str(v) if not isinstance(v,(list,dict)) else json.dumps(v))
            r += 1
        for sec in (data.get('sections') or []):
            ws.cell(row=r, column=1, value=sec.get('heading','')).font = Font(bold=True, color='1B5E20')
            ws.cell(row=r, column=2, value=sec.get('body',''))
            r += 1

    for col_letter, width in [('A', 28), ('B', 48), ('C', 14), ('D', 14)]:
        ws.column_dimensions[col_letter].width = width
    wb.save(out)


def _build_pptx(doc_type, data, out):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    p = Presentation()
    title = data.get('title') or DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document')

    # Title slide
    s = p.slides.add_slide(p.slide_layouts[0])
    s.shapes.title.text = title
    if s.placeholders and len(s.placeholders) > 1:
        s.placeholders[1].text = data.get('subtitle','Generated by Briefly')

    sections = data.get('sections') or []
    if not sections and data.get('content'):
        content = _expand_with_llm(doc_type, data)
        # Split into chunks of ~5 lines for each slide
        paras = [p.strip() for p in content.split('\n\n') if p.strip()]
        sections = [{'heading': f'Slide {i+1}', 'body': para} for i, para in enumerate(paras)]

    for sec in sections:
        sl = p.slides.add_slide(p.slide_layouts[1])
        sl.shapes.title.text = sec.get('heading','')
        body = sl.placeholders[1].text_frame
        body.text = ''
        for line in str(sec.get('body','')).split('\n'):
            if not line.strip(): continue
            para = body.add_paragraph()
            para.text = line.strip()
            para.font.size = Pt(18)
    p.save(out)


def _build_txt(doc_type, data, out):
    title = data.get('title') or DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document')
    lines = [title, '=' * len(title), '']
    if doc_type == 'bill':
        lines += [
            f"Invoice #: {data.get('invoice_number','INV-001')}",
            f"Date: {data.get('date', dt.date.today().isoformat())}",
            f"Bill To: {data.get('client_name','')}",
            '',
            f"{'Item':<30}{'Qty':>6}{'Price':>10}{'Total':>10}",
            '-' * 56,
        ]
        total = 0.0
        for it in (data.get('items') or []):
            qty = float(it.get('qty',1)); price = float(it.get('price',0)); line = qty*price; total += line
            lines.append(f"{str(it.get('name',''))[:28]:<30}{qty:>6g}{price:>10.2f}{line:>10.2f}")
        lines += ['-'*56, f"{'Total':>46}{data.get('currency','$')}{total:.2f}"]
    elif doc_type == 'resume':
        if data.get('full_name'): lines.append(data['full_name'])
        if data.get('headline'): lines.append(data['headline'])
        contact = ' • '.join(filter(None, [data.get('email'), data.get('phone'), data.get('location'), data.get('website')]))
        if contact: lines.append(contact)
        for k, label in [('summary','PROFILE'), ('experience','WORK EXPERIENCE'),
                         ('education','EDUCATION'), ('skills','TECHNICAL SKILLS'),
                         ('projects','PROJECTS'), ('hobbies','HOBBIES'),
                         ('languages','LANGUAGES KNOWN')]:
            if data.get(k):
                lines += ['', label, '-'*len(label), str(data[k])]
    else:
        for sec in (data.get('sections') or []):
            lines += ['', sec.get('heading',''), '-'*len(sec.get('heading','')), sec.get('body','')]
        if data.get('content'):
            lines += ['', _expand_with_llm(doc_type, data)]
    with open(out, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))


def _build_md(doc_type, data, out):
    title = data.get('title') or DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document')
    lines = [f"# {title}", '']
    if doc_type == 'bill':
        lines += [
            f"**Invoice #:** {data.get('invoice_number','INV-001')}  ",
            f"**Date:** {data.get('date', dt.date.today().isoformat())}  ",
            f"**Bill To:** {data.get('client_name','')}", '',
            '| Item | Qty | Price | Total |', '|---|---:|---:|---:|',
        ]
        total = 0.0
        for it in (data.get('items') or []):
            qty = float(it.get('qty',1)); price = float(it.get('price',0)); line = qty*price; total += line
            lines.append(f"| {it.get('name','')} | {qty:g} | {price:.2f} | {line:.2f} |")
        lines += ['', f"**Total: {data.get('currency','$')}{total:.2f}**"]
    elif doc_type == 'resume':
        if data.get('full_name'): lines.append(f"## {data['full_name']}")
        if data.get('headline'): lines.append(f"*{data['headline']}*")
        contact = ' • '.join(filter(None, [data.get('email'), data.get('phone'), data.get('location'), data.get('website')]))
        if contact: lines.append(contact)
        for k, label in [('summary','Profile'), ('experience','Work Experience'),
                         ('education','Education'), ('skills','Technical Skills'),
                         ('projects','Projects'), ('hobbies','Hobbies'),
                         ('languages','Languages Known')]:
            if data.get(k):
                lines += ['', f'### {label}', str(data[k])]
    else:
        for sec in (data.get('sections') or []):
            lines += ['', f"## {sec.get('heading','')}", sec.get('body','')]
        if data.get('content'):
            lines += ['', _expand_with_llm(doc_type, data)]
    with open(out, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))


BUILDERS = {
    'docx': (_build_docx, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'),
    'pdf':  (_build_pdf,  'application/pdf'),
    'xlsx': (_build_xlsx, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'),
    'pptx': (_build_pptx, 'application/vnd.openxmlformats-officedocument.presentationml.presentation'),
    'txt':  (_build_txt,  'text/plain'),
    'md':   (_build_md,   'text/markdown'),
}


@app.route('/api/document-maker/generate', methods=['POST'])
def doc_maker_generate():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    payload = request.json or {}
    doc_type = payload.get('doc_type', 'custom')
    fmt = (payload.get('format') or 'pdf').lower()
    data = payload.get('data') or {}
    if fmt not in BUILDERS:
        return jsonify({'success': False, 'error': f'Unsupported format: {fmt}'}), 400
    if doc_type not in DOC_TEMPLATES:
        return jsonify({'success': False, 'error': 'Unknown document type'}), 400

    builder, mime = BUILDERS[fmt]
    base = _safe_name(data.get('title') or DOC_TEMPLATES[doc_type]['label'])
    out_name = f"{base}.{fmt}"
    out_path = os.path.join(app.config['UPLOAD_FOLDER'], out_name)
    try:
        builder(doc_type, data, out_path)
    except Exception as e:
        return jsonify({'success': False, 'error': f'Generation failed: {e}'}), 500

    return send_file(out_path, mimetype=mime, as_attachment=True, download_name=out_name)


# Root index — small JSON to confirm API is up (HTML is served from frontend)
@app.route('/api/health')
def health():
    return jsonify({'success': True, 'status': 'ok', 'service': 'Briefly API'})


# =========== STREAMING CHAT (SSE) ===========
@app.route('/api/chat/stream', methods=['POST'])
def chat_stream():
    data = request.json or {}
    user_message = (data.get('message') or '').strip()
    conversation_id = data.get('conversation_id')
    if not user_message:
        return jsonify({'success': False, 'error': 'Message cannot be empty'}), 400

    # Auto-create conversation
    if not conversation_id:
        user_id = session.get('user_id', 0)
        title = user_message[:50] + ('...' if len(user_message) > 50 else '')
        conversation_id = create_conversation(user_id, title)

    if 'user_id' in session:
        conn = get_db(); cur = conn.cursor()
        cur.execute('SELECT user_id FROM conversations WHERE id = ?', (conversation_id,))
        row = cur.fetchone(); conn.close()
        if row and row[0] != session['user_id'] and row[0] != 0:
            return jsonify({'success': False, 'error': 'Access denied'}), 403

    # Save user message immediately so refresh keeps it
    save_message(conversation_id, 'user', user_message)

    lower = user_message.lower().strip()
    simple = simple_responses.get(lower)

    def event_stream():
        # First event: meta (conversation_id + user message timestamp)
        meta = {'conversation_id': conversation_id, 'created_at': dt.datetime.utcnow().isoformat() + 'Z'}
        yield f"event: meta\ndata: {json.dumps(meta)}\n\n"

        full_reply = ''
        try:
            if simple:
                # Stream the simple response char-by-char for a snappy feel
                full_reply = simple
                for ch in simple:
                    yield f"data: {json.dumps({'delta': ch})}\n\n"
            else:
                # Build prompt with history
                previous = get_conversation_messages(conversation_id)
                # Drop the user message we just saved (it's the last one)
                history_lines = []
                for m in previous[-10:-1]:
                    tag = 'User' if m['role'] == 'user' else 'Assistant'
                    history_lines.append(f"{tag}: {m['content']}")
                history_text = '\n'.join(history_lines)
                prompt = (
                    "You are Briefly, a helpful, concise and friendly AI assistant. "
                    "Use markdown formatting when useful (headings, lists, **bold**, `code`).\n\n"
                )
                if history_text:
                    prompt += f"Conversation so far:\n{history_text}\n\n"
                prompt += f"User: {user_message}\nAssistant:"

                # Use Gemini streaming
                stream = client.models.generate_content_stream(
                    model=GEMINI_MODEL,
                    contents=prompt,
                )
                for chunk in stream:
                    text = getattr(chunk, 'text', None)
                    if text:
                        full_reply += text
                        yield f"data: {json.dumps({'delta': text})}\n\n"
        except Exception as e:
            yield f"event: error\ndata: {json.dumps({'error': str(e)})}\n\n"
            return

        # Save assistant message
        if full_reply:
            save_message(conversation_id, 'assistant', full_reply)
        yield f"event: done\ndata: {json.dumps({'full': full_reply, 'created_at': dt.datetime.utcnow().isoformat() + 'Z'})}\n\n"

    return Response(stream_with_context(event_stream()),
                    mimetype='text/event-stream',
                    headers={
                        'Cache-Control': 'no-cache, no-transform',
                        'X-Accel-Buffering': 'no',
                        'Connection': 'keep-alive',
                    })


# =========== DOCUMENT HISTORY (Analyzer) ===========
@app.route('/api/documents', methods=['GET'])
def list_documents():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    docs = list_analyzed_documents(session['user_id'])
    return jsonify({'success': True, 'documents': docs})


@app.route('/api/documents/<int:doc_id>', methods=['GET'])
def fetch_document(doc_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    doc = get_analyzed_document(session['user_id'], doc_id)
    if not doc:
        return jsonify({'success': False, 'error': 'Not found'}), 404
    return jsonify({'success': True, 'document': doc})


@app.route('/api/documents/<int:doc_id>/delete', methods=['POST'])
def remove_document(doc_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    if delete_analyzed_document(session['user_id'], doc_id):
        return jsonify({'success': True})
    return jsonify({'success': False, 'error': 'Not found'}), 404


# =========== EDITOR DRAFTS (Scribe text editor & Gridly spreadsheet) ===========
@app.route('/api/editor/drafts', methods=['GET'])
def api_drafts_list():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    kind = request.args.get('kind')
    return jsonify({'success': True, 'drafts': list_drafts(session['user_id'], kind)})


@app.route('/api/editor/drafts', methods=['POST'])
def api_drafts_create():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    d = request.json or {}
    kind = d.get('kind')
    if kind not in ('text', 'sheet'):
        return jsonify({'success': False, 'error': 'kind must be text or sheet'}), 400
    title = (d.get('title') or 'Untitled').strip()
    content = d.get('content') or ''
    did = create_draft(session['user_id'], kind, title, content)
    return jsonify({'success': True, 'id': did})


@app.route('/api/editor/drafts/<int:draft_id>', methods=['GET'])
def api_drafts_get(draft_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    draft = get_draft(session['user_id'], draft_id)
    if not draft:
        return jsonify({'success': False, 'error': 'Not found'}), 404
    return jsonify({'success': True, 'draft': draft})


@app.route('/api/editor/drafts/<int:draft_id>', methods=['PUT'])
def api_drafts_update(draft_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    d = request.json or {}
    ok = update_draft(session['user_id'], draft_id, d.get('title'), d.get('content'))
    if not ok: return jsonify({'success': False, 'error': 'Not found'}), 404
    return jsonify({'success': True})


@app.route('/api/editor/drafts/<int:draft_id>/delete', methods=['POST'])
def api_drafts_delete(draft_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    if delete_draft(session['user_id'], draft_id):
        return jsonify({'success': True})
    return jsonify({'success': False, 'error': 'Not found'}), 404


# ---- Scribe (text) export ----
def _html_to_docx_bytes(html, title='Document'):
    from io import BytesIO
    from bs4 import BeautifulSoup
    from docx import Document
    from docx.shared import Pt, RGBColor

    doc = Document()
    soup = BeautifulSoup(html or '', 'html.parser')
    body = soup.body or soup

    def add_runs(p, node, bold=False, italic=False, underline=False):
        if hasattr(node, 'name') and node.name is None:
            return
        if not hasattr(node, 'name') or node.name is None:
            # text node
            txt = str(node)
            if txt:
                r = p.add_run(txt)
                r.bold = bold; r.italic = italic; r.underline = underline
                r.font.color.rgb = RGBColor(0, 0, 0)
            return
        name = node.name.lower() if node.name else ''
        new_bold = bold or name in ('b', 'strong')
        new_italic = italic or name in ('i', 'em')
        new_underline = underline or name == 'u'
        if name == 'br':
            p.add_run().add_break(); return
        if name == 'a':
            text = node.get_text()
            url = node.get('href') or text
            try:
                _add_hyperlink(p, url, text, bold=new_bold)
                return
            except Exception:
                pass
        for child in node.children:
            if hasattr(child, 'children') and getattr(child, 'name', None):
                add_runs(p, child, new_bold, new_italic, new_underline)
            else:
                txt = str(child)
                if txt:
                    r = p.add_run(txt)
                    r.bold = new_bold; r.italic = new_italic; r.underline = new_underline
                    r.font.color.rgb = RGBColor(0, 0, 0)

    def walk(node):
        for child in node.children:
            if not hasattr(child, 'name') or child.name is None:
                txt = str(child).strip()
                if txt:
                    p = doc.add_paragraph(); add_runs(p, child)
                continue
            tag = child.name.lower()
            if tag in ('h1', 'h2', 'h3', 'h4', 'h5', 'h6'):
                level = int(tag[1])
                p = doc.add_heading('', level=min(level, 4))
                for c in child.children: add_runs(p, c)
                for r in p.runs: r.font.color.rgb = RGBColor(0, 0, 0)
            elif tag == 'p' or tag == 'div':
                p = doc.add_paragraph()
                for c in child.children: add_runs(p, c)
            elif tag in ('ul', 'ol'):
                style = 'List Bullet' if tag == 'ul' else 'List Number'
                for li in child.find_all('li', recursive=False):
                    p = doc.add_paragraph(style=style)
                    for c in li.children: add_runs(p, c)
            elif tag == 'blockquote':
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Pt(24)
                for c in child.children: add_runs(p, c, italic=True)
            elif tag == 'pre':
                p = doc.add_paragraph()
                r = p.add_run(child.get_text())
                r.font.name = 'Consolas'
            elif tag == 'hr':
                doc.add_paragraph('—' * 30)
            elif tag == 'table':
                rows = child.find_all('tr')
                if not rows: continue
                cols = max(len(r.find_all(['td', 'th'])) for r in rows)
                tbl = doc.add_table(rows=len(rows), cols=cols)
                tbl.style = 'Table Grid'
                for ri, tr in enumerate(rows):
                    for ci, cell in enumerate(tr.find_all(['td', 'th'])):
                        tbl.rows[ri].cells[ci].text = cell.get_text()
            else:
                # fallback: dump as paragraph
                p = doc.add_paragraph()
                for c in child.children: add_runs(p, c)

    walk(body)
    buf = BytesIO(); doc.save(buf); buf.seek(0); return buf


def _html_to_pdf_bytes(html, title='Document'):
    from io import BytesIO
    from bs4 import BeautifulSoup
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable

    buf = BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=1.8*cm, rightMargin=1.8*cm, topMargin=2*cm, bottomMargin=2*cm)
    ss = getSampleStyleSheet()
    BLACK = colors.HexColor('#000000')
    body = ParagraphStyle('rb', parent=ss['BodyText'], textColor=BLACK, fontSize=11, leading=15, spaceAfter=6)
    h1 = ParagraphStyle('rh1', parent=ss['Heading1'], textColor=BLACK, fontSize=18, spaceAfter=8)
    h2 = ParagraphStyle('rh2', parent=ss['Heading2'], textColor=BLACK, fontSize=15, spaceAfter=6)
    h3 = ParagraphStyle('rh3', parent=ss['Heading3'], textColor=BLACK, fontSize=13, spaceAfter=4)
    soup = BeautifulSoup(html or '', 'html.parser')
    body_node = soup.body or soup
    story = []

    def inline_html(node):
        # Convert child HTML to reportlab-compatible inline tags (b/i/u/a/font)
        from copy import copy
        s = ''
        for c in node.children:
            if not hasattr(c, 'name') or c.name is None:
                s += str(c).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            else:
                tag = c.name.lower()
                inner = inline_html(c)
                if tag in ('b', 'strong'): s += f'<b>{inner}</b>'
                elif tag in ('i', 'em'): s += f'<i>{inner}</i>'
                elif tag == 'u': s += f'<u>{inner}</u>'
                elif tag == 'a':
                    href = c.get('href') or inner
                    s += f'<a href="{href}"><font color="#0563C1"><u>{inner}</u></font></a>'
                elif tag == 'br': s += '<br/>'
                else: s += inner
        return s

    for child in body_node.children:
        if not hasattr(child, 'name') or child.name is None:
            txt = str(child).strip()
            if txt: story.append(Paragraph(txt, body))
            continue
        tag = child.name.lower()
        if tag in ('h1', 'h2', 'h3', 'h4'):
            style = {'h1': h1, 'h2': h2}.get(tag, h3)
            story.append(Paragraph(inline_html(child), style))
        elif tag == 'p' or tag == 'div':
            story.append(Paragraph(inline_html(child), body))
        elif tag == 'hr':
            story.append(HRFlowable(width='100%', thickness=0.5, color=BLACK, spaceBefore=4, spaceAfter=6))
        elif tag in ('ul', 'ol'):
            for i, li in enumerate(child.find_all('li', recursive=False), 1):
                prefix = '• ' if tag == 'ul' else f'{i}. '
                story.append(Paragraph(prefix + inline_html(li), body))
        elif tag == 'blockquote':
            story.append(Paragraph('<i>' + inline_html(child) + '</i>', body))
        elif tag == 'pre':
            story.append(Paragraph('<font name="Courier">' + child.get_text().replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br/>') + '</font>', body))
        elif tag == 'table':
            rows = child.find_all('tr')
            data_rows = [[cell.get_text() for cell in tr.find_all(['td', 'th'])] for tr in rows]
            if data_rows:
                t = Table(data_rows)
                t.setStyle(TableStyle([('GRID', (0,0), (-1,-1), 0.4, colors.grey)]))
                story.append(t); story.append(Spacer(1, 6))
        else:
            story.append(Paragraph(inline_html(child), body))
    if not story: story.append(Paragraph(' ', body))
    doc.build(story)
    buf.seek(0); return buf


@app.route('/api/editor/export/text', methods=['POST'])
def api_editor_export_text():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    d = request.json or {}
    fmt = (d.get('format') or 'pdf').lower()
    title = _safe_name(d.get('title') or 'document')
    html = d.get('html') or ''
    try:
        if fmt == 'docx':
            buf = _html_to_docx_bytes(html, title)
            return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                             as_attachment=True, download_name=f'{title}.docx')
        if fmt == 'doc':
            # python-docx writes .docx; alias .doc to .docx file for compatibility
            buf = _html_to_docx_bytes(html, title)
            return send_file(buf, mimetype='application/msword',
                             as_attachment=True, download_name=f'{title}.doc')
        if fmt == 'pdf':
            buf = _html_to_pdf_bytes(html, title)
            return send_file(buf, mimetype='application/pdf',
                             as_attachment=True, download_name=f'{title}.pdf')
        if fmt == 'txt':
            from bs4 import BeautifulSoup
            text = BeautifulSoup(html, 'html.parser').get_text(separator='\n')
            from io import BytesIO
            buf = BytesIO(text.encode('utf-8')); buf.seek(0)
            return send_file(buf, mimetype='text/plain', as_attachment=True, download_name=f'{title}.txt')
        return jsonify({'success': False, 'error': f'Unsupported format: {fmt}'}), 400
    except Exception as e:
        return jsonify({'success': False, 'error': f'Export failed: {e}'}), 500


@app.route('/api/editor/export/sheet', methods=['POST'])
def api_editor_export_sheet():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401
    payload = request.json or {}
    fmt = (payload.get('format') or 'xlsx').lower()
    title = _safe_name(payload.get('title') or 'spreadsheet')
    rows = payload.get('rows') or []  # 2D list of strings/numbers
    try:
        from io import BytesIO
        if fmt == 'csv':
            import csv, io
            s = io.StringIO()
            w = csv.writer(s)
            for row in rows:
                w.writerow(row if isinstance(row, list) else [row])
            data = s.getvalue().encode('utf-8')
            buf = BytesIO(data); buf.seek(0)
            return send_file(buf, mimetype='text/csv', as_attachment=True, download_name=f'{title}.csv')
        if fmt == 'xlsx':
            from openpyxl import Workbook
            wb = Workbook(); ws = wb.active; ws.title = title[:30] or 'Sheet1'
            for ri, row in enumerate(rows, 1):
                for ci, val in enumerate(row, 1):
                    # try numeric conversion
                    try:
                        v = float(val)
                        if v.is_integer(): v = int(v)
                    except (TypeError, ValueError):
                        v = val
                    ws.cell(row=ri, column=ci, value=v)
            buf = BytesIO(); wb.save(buf); buf.seek(0)
            return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                             as_attachment=True, download_name=f'{title}.xlsx')
        return jsonify({'success': False, 'error': f'Unsupported format: {fmt}'}), 400
    except Exception as e:
        return jsonify({'success': False, 'error': f'Export failed: {e}'}), 500
    
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(
        host="0.0.0.0",
        port=port,
        debug=False
    )