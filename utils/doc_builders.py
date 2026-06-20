import re
import datetime as dt
from flask import current_app
import docx
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

# Structuring Constants
DOC_TEMPLATES = {
    'bill':    {'label': 'Bill / Invoice',          'icon': 'fa-file-invoice-dollar'},
    'resume':  {'label': 'Resume / CV',             'icon': 'fa-id-card'},
    'project_architecture': {'label': 'Project Architecture', 'icon': 'fa-diagram-project'},
    'school_project':  {'label': 'School Project',  'icon': 'fa-school'},
    'college_project': {'label': 'College Project', 'icon': 'fa-graduation-cap'},
    'thesis':  {'label': 'Thesis / Dissertation',   'icon': 'fa-book'},
    'report':  {'label': 'Report',                  'icon': 'fa-file-lines'},
    'letter':  {'label': 'Letter',                  'icon': 'fa-envelope-open-text'},
    'meeting_minutes': {'label': 'Meeting Minutes', 'icon': 'fa-users'},
    'business_plan':  {'label': 'Business Plan',    'icon': 'fa-briefcase'},
    'custom':  {'label': 'Custom Document',         'icon': 'fa-file-pen'},
}

RESUME_TEMPLATES = ['classic', 'modern', 'compact']

def _expand_with_llm(doc_type, data_dict):
    raw = data_dict.get('content', '') or ''
    if not raw or len(raw) > 200:
        return raw  
    try:
        client = current_app.config['OPENAI_CLIENT']
        model_name = current_app.config['GEMINI_MODEL']
        
        prompt = (
            f"Expand the following short prompt into a well-structured {DOC_TEMPLATES.get(doc_type,{}).get('label','document')} "
            f"of about 300-500 words. Return plain text with clear paragraph breaks; do NOT add markdown.\n\nPrompt: {raw}"
        )
        resp = client.chat.completions.create(
            model=model_name,
            messages=[{"role": "user", "content": prompt}]
        )
        return (resp.choices[0].message.content or raw).strip()
    except Exception:
        return raw

# Docx Helper Logic
def _add_hyperlink(paragraph, url, text, bold=False, size=None):
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    part = paragraph.part
    r_id = part.relate_to(url, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
    hyperlink = OxmlElement('w:hyperlink')
    hyperlink.set(qn('r:id'), r_id)
    new_run = OxmlElement('w:r')
    rPr = OxmlElement('w:rPr')
    color = OxmlElement('w:color'); color.set(qn('w:val'), '0563C1'); rPr.append(color)
    u = OxmlElement('w:u'); u.set(qn('w:val'), 'single'); rPr.append(u)
    if bold:
        b = OxmlElement('w:b'); rPr.append(b)
    if size:
        sz = OxmlElement('w:sz'); sz.set(qn('w:val'), str(size * 2)); rPr.append(sz)
    new_run.append(rPr)
    txt = OxmlElement('w:t'); txt.text = text; txt.set(qn('xml:space'), 'preserve')
    new_run.append(txt)
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)

def _docx_inline_contacts(p, text, size=11):
    pattern = re.compile(r'(https?://[^\s]+|[\w._%+-]+@[\w.-]+\.[A-Za-z]{2,}|\+?\d[\d\s().-]{6,}\d)')
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            run = p.add_run(text[pos:m.start()])
            run.font.color.rgb = RGBColor(0, 0, 0)
            if size: run.font.size = Pt(size)
        target = m.group(0)
        if '@' in target and not target.startswith('http'):
            href = 'mailto:' + target
        elif (target.startswith('+') or target[0].isdigit()) and not target.startswith('http'):
            href = 'tel:' + re.sub(r'[\s().-]', '', target)
        else:
            href = target
        _add_hyperlink(p, href, target, size=size)
        pos = m.end()
    if pos < len(text):
        run = p.add_run(text[pos:])
        run.font.color.rgb = RGBColor(0, 0, 0)
        if size: run.font.size = Pt(size)

def _build_resume_docx(d, data):
    template = (data.get('template') or 'classic').lower()
    if template not in RESUME_TEMPLATES: template = 'classic'

    name = data.get('full_name', '').strip()
    title = data.get('title', '').strip() if data.get('title') and data.get('title') != 'Resume' else (data.get('headline') or '').strip()
    contact_bits = [b for b in [data.get('email'), data.get('phone'), data.get('location'), data.get('website')] if b]

    if template == 'modern':
        p = d.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if name:
            r = p.add_run(name); r.bold = True; r.font.size = Pt(20); r.font.color.rgb = RGBColor(0, 0, 0)
        if title:
            p2 = d.add_paragraph(); p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r = p2.add_run(title); r.font.size = Pt(11); r.font.color.rgb = RGBColor(0, 0, 0)
        if contact_bits:
            p3 = d.add_paragraph(); p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _docx_inline_contacts(p3, '  |  '.join(contact_bits), size=10)
    else:
        if name:
            p = d.add_paragraph()
            r = p.add_run(name); r.bold = True
            r.font.size = Pt(18 if template == 'classic' else 16)
            r.font.color.rgb = RGBColor(0, 0, 0)
        if title:
            p = d.add_paragraph()
            r = p.add_run(title); r.font.size = Pt(11); r.font.color.rgb = RGBColor(0, 0, 0)
        if contact_bits:
            p = d.add_paragraph()
            _docx_inline_contacts(p, ' • '.join(contact_bits), size=10)

    section_size = 12 if template == 'compact' else 13

    def section(heading):
        p = d.add_paragraph()
        r = p.add_run(heading.upper() if template == 'compact' else heading)
        r.bold = True; r.font.size = Pt(section_size); r.font.color.rgb = RGBColor(0, 0, 0)

    def body_para(text, bullet=False):
        p = d.add_paragraph(style='List Bullet') if bullet else d.add_paragraph()
        _docx_inline_contacts(p, text, size=10 if template == 'compact' else 11)

    if data.get('summary'):
        section('Profile'); body_para(data['summary'])
    if data.get('experience'):
        section('Work Experience')
        for line in str(data['experience']).split('\n'):
            if line.strip(): body_para(line.strip(), bullet=(template != 'classic'))
    if data.get('education'):
        section('Education')
        for line in str(data['education']).split('\n'):
            if line.strip(): body_para(line.strip(), bullet=(template != 'classic'))
    if data.get('skills'):
        section('Technical Skills'); body_para(data['skills'])
    if data.get('projects'):
        section('Projects')
        for line in str(data['projects']).split('\n'):
            if line.strip(): body_para(line.strip(), bullet=(template != 'classic'))
    if data.get('hobbies'):
        section('Hobbies'); body_para(data['hobbies'])
    if data.get('languages'):
        section('Languages Known'); body_para(data['languages'])

def _build_docx(doc_type, data, out):
    d = docx.Document()
    if doc_type != 'resume':
        title = data.get('title') or DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document')
        h = d.add_heading(title, level=0)
        for run in h.runs: run.font.color.rgb = RGBColor(0x1B, 0x5E, 0x20)

    if doc_type == 'resume':
        _build_resume_docx(d, data)
    elif doc_type == 'bill':
        d.add_paragraph(f"Invoice #: {data.get('invoice_number','INV-001')}")
        d.add_paragraph(f"Date: {data.get('date', dt.date.today().isoformat())}")
        d.add_paragraph(f"Bill To: {data.get('client_name','')}")
        if data.get('client_address'): d.add_paragraph(data['client_address'])
        d.add_paragraph(' ')
        items = data.get('items') or []
        tbl = d.add_table(rows=1, cols=4); tbl.style = 'Light Grid Accent 1'
        hdr = tbl.rows[0].cells
        hdr[0].text, hdr[1].text, hdr[2].text, hdr[3].text = 'Item', 'Qty', 'Price', 'Total'
        total = 0.0
        for it in items:
            row = tbl.add_row().cells
            qty = float(it.get('qty', 1)); price = float(it.get('price', 0))
            line = qty * price; total += line
            row[0].text = str(it.get('name', ''))
            row[1].text = str(qty); row[2].text = f"{price:.2f}"; row[3].text = f"{line:.2f}"
        d.add_paragraph(' ')
        p = d.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        run = p.add_run(f"Total: {data.get('currency','$')}{total:.2f}"); run.bold = True
        if data.get('notes'): d.add_paragraph(data['notes'])
    else:
        content = _expand_with_llm(doc_type, data) if data.get('content') else ''
        sections = data.get('sections') or []
        for sec in sections:
            if sec.get('heading'): d.add_heading(sec['heading'], level=1)
            if sec.get('body'): d.add_paragraph(sec['body'])
        if content:
            for para in content.split('\n\n'):
                if para.strip(): d.add_paragraph(para.strip())
    d.save(out)

# ReportLab PDF Helper Logic
def _build_resume_pdf(story, data, ss):
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib import colors
    from reportlab.platypus import Paragraph, HRFlowable

    template = (data.get('template') or 'classic').lower()
    if template not in RESUME_TEMPLATES: template = 'classic'

    BLACK = colors.HexColor('#000000')
    name_style = ParagraphStyle('rname', parent=ss['Title'], textColor=BLACK, fontSize=22 if template != 'compact' else 18, alignment=1 if template == 'modern' else 0, spaceAfter=2, leading=24)
    title_style2 = ParagraphStyle('rtitle', parent=ss['BodyText'], textColor=BLACK, fontSize=11, alignment=1 if template == 'modern' else 0, spaceAfter=6)
    contact_style = ParagraphStyle('rcontact', parent=ss['BodyText'], textColor=BLACK, fontSize=10, alignment=1 if template == 'modern' else 0, spaceAfter=10)
    section_style = ParagraphStyle('rsection', parent=ss['BodyText'], textColor=BLACK, fontSize=12 if template == 'compact' else 13, fontName='Helvetica-Bold', spaceBefore=10, spaceAfter=4)
    body_style = ParagraphStyle('rbody', parent=ss['BodyText'], textColor=BLACK, fontSize=10 if template == 'compact' else 11, spaceAfter=4, leading=14)

    def linkify(text):
        def repl(m):
            t = m.group(0)
            href = 'mailto:' + t if '@' in t and not t.startswith('http') else ('tel:' + re.sub(r'[\s().-]', '', t) if (t.startswith('+') or t[0].isdigit()) and not t.startswith('http') else t)
            return f'<font color="#0563C1"><u><a href="{href}">{t}</a></u></font>'
        return re.sub(r'https?://[^\s<]+|[\w._%+-]+@[\w.-]+\.[A-Za-z]{2,}|\+?\d[\d\s().-]{6,}\d', repl, text)

    name = data.get('full_name', '').strip()
    headline = (data.get('headline') or '').strip()
    contact_bits = [b for b in [data.get('email'), data.get('phone'), data.get('location'), data.get('website')] if b]

    if name: story.append(Paragraph(f"<b>{name}</b>", name_style))
    if headline: story.append(Paragraph(headline, title_style2))
    if contact_bits: story.append(Paragraph(linkify(('  |  ' if template == 'modern' else ' &bull; ').join(contact_bits)), contact_style))
    if template == 'modern': story.append(HRFlowable(width='100%', thickness=0.5, color=BLACK, spaceBefore=2, spaceAfter=8))

    def add_section(title, value, bulletize=True):
        if not value: return
        story.append(Paragraph(title.upper() if template == 'compact' else title, section_style))
        if isinstance(value, str) and '\n' in value and bulletize:
            for line in value.split('\n'):
                if line.strip(): story.append(Paragraph(linkify(line.strip()) if template == 'classic' else '• ' + linkify(line.strip()), body_style))
        else:
            story.append(Paragraph(linkify(str(value)), body_style))

    add_section('Profile', data.get('summary'), bulletize=False)
    add_section('Work Experience', data.get('experience'))
    add_section('Education', data.get('education'))
    add_section('Technical Skills', data.get('skills'), bulletize=False)
    add_section('Projects', data.get('projects'))
    add_section('Hobbies', data.get('hobbies'), bulletize=False)
    add_section('Languages Known', data.get('languages'), bulletize=False)

def _build_pdf(doc_type, data, out):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle

    doc = SimpleDocTemplate(out, pagesize=A4, leftMargin=1.8*cm, rightMargin=1.8*cm, topMargin=2*cm, bottomMargin=2*cm)
    ss = getSampleStyleSheet()
    title_style = ParagraphStyle('t', parent=ss['Title'], textColor=colors.HexColor('#1B5E20'))
    body = ss['BodyText']; body.spaceAfter = 6
    h1 = ParagraphStyle('h1', parent=ss['Heading2'], textColor=colors.HexColor('#1B5E20'))
    story = []
    
    if doc_type != 'resume':
        title = data.get('title') or DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document')
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 12))

    if doc_type == 'bill':
        story.append(Paragraph(f"Invoice #: {data.get('invoice_number','INV-001')}", body))
        story.append(Paragraph(f"Date: {data.get('date', dt.date.today().isoformat())}", body))
        story.append(Paragraph(f"Bill To: {data.get('client_name','')}", body))
        if data.get('client_address'): story.append(Paragraph(data['client_address'], body))
        story.append(Spacer(1, 10))
        rows = [['Item', 'Qty', 'Price', 'Total']]; total = 0.0
        for it in (data.get('items') or []):
            qty = float(it.get('qty', 1)); price = float(it.get('price', 0))
            line = qty * price; total += line
            rows.append([it.get('name',''), f"{qty:g}", f"{price:.2f}", f"{line:.2f}"])
        rows.append(['', '', 'Total', f"{data.get('currency','$')}{total:.2f}"])
        t = Table(rows, colWidths=[7*cm, 2*cm, 3*cm, 3*cm])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1B5E20')),
            ('TEXTCOLOR', (0,0), (-1,0), colors.white),
            ('GRID', (0,0), (-1,-1), 0.4, colors.grey),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('BACKGROUND', (0,-1), (-1,-1), colors.HexColor('#E8F5E9')),
        ]))
        story.append(t)
        if data.get('notes'):
            story.append(Spacer(1, 12)); story.append(Paragraph(data['notes'], body))
    elif doc_type == 'resume':
        _build_resume_pdf(story, data, ss)
    else:
        content = _expand_with_llm(doc_type, data) if data.get('content') else ''
        for sec in (data.get('sections') or []):
            if sec.get('heading'): story.append(Paragraph(sec['heading'], h1))
            if sec.get('body'): story.append(Paragraph(sec['body'].replace('\n','<br/>'), body))
        if content:
            for para in content.split('\n\n'):
                if para.strip(): story.append(Paragraph(para.strip().replace('\n','<br/>'), body))
    doc.build(story)

# Additional formats
def _build_xlsx(doc_type, data, out):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    wb = Workbook(); ws = wb.active
    sheet_title = re.sub(r'[\\/?*\[\]:]', '-', DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document'))
    ws.title = sheet_title[:30]
    title = data.get('title') or DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document')
    ws['A1'] = title; ws['A1'].font = Font(size=16, bold=True, color='1B5E20')
    ws.merge_cells('A1:D1'); ws['A1'].alignment = Alignment(horizontal='center')

    if doc_type == 'bill':
        ws['A3'] = 'Invoice #'; ws['B3'] = data.get('invoice_number','INV-001')
        ws['A4'] = 'Date';      ws['B4'] = data.get('date', dt.date.today().isoformat())
        ws['A5'] = 'Bill To';   ws['B5'] = data.get('client_name','')
        headers = ['Item','Qty','Price','Total']
        for i, h in enumerate(headers, 1):
            c = ws.cell(row=7, column=i, value=h)
            c.font = Font(bold=True, color='FFFFFF'); c.fill = PatternFill('solid', fgColor='1B5E20')
        r = 8; total = 0.0
        for it in (data.get('items') or []):
            qty = float(it.get('qty',1)); price = float(it.get('price',0)); line = qty*price; total += line
            ws.cell(row=r, column=1, value=it.get('name',''))
            ws.cell(row=r, column=2, value=qty); ws.cell(row=r, column=3, value=price); ws.cell(row=r, column=4, value=line)
            r += 1
        ws.cell(row=r, column=3, value='Total').font = Font(bold=True)
        ws.cell(row=r, column=4, value=total).font = Font(bold=True)
    else:
        r = 3
        for k, v in data.items():
            if k in ('title','sections') or v is None: continue
            ws.cell(row=r, column=1, value=str(k).replace('_',' ').title()).font = Font(bold=True)
            ws.cell(row=r, column=2, value=str(v) if not isinstance(v,(list,dict)) else json.dumps(v))
            r += 1
        for sec in (data.get('sections') or []):
            ws.cell(row=r, column=1, value=sec.get('heading','')).font = Font(bold=True, color='1B5E20')
            ws.cell(row=r, column=2, value=sec.get('body',''))
            r += 1
    for col_letter, width in [('A', 28), ('B', 48), ('C', 14), ('D', 14)]: ws.column_dimensions[col_letter].width = width
    wb.save(out)

def _build_pptx(doc_type, data, out):
    from pptx import Presentation
    from pptx.util import Pt
    p = Presentation()
    title = data.get('title') or DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document')
    s = p.slides.add_slide(p.slide_layouts[0])
    s.shapes.title.text = title
    if s.placeholders and len(s.placeholders) > 1: s.placeholders[1].text = data.get('subtitle','Generated by Briefly')
    sections = data.get('sections') or []
    if not sections and data.get('content'):
        paras = [p.strip() for p in _expand_with_llm(doc_type, data).split('\n\n') if p.strip()]
        sections = [{'heading': f'Slide {i+1}', 'body': para} for i, para in enumerate(paras)]
    for sec in sections:
        sl = p.slides.add_slide(p.slide_layouts[1])
        sl.shapes.title.text = sec.get('heading','')
        body = sl.placeholders[1].text_frame; body.text = ''
        for line in str(sec.get('body','')).split('\n'):
            if line.strip():
                para = body.add_paragraph(); para.text = line.strip(); para.font.size = Pt(18)
    p.save(out)

def _build_txt(doc_type, data, out):
    title = data.get('title') or DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document')
    lines = [title, '=' * len(title), '']
    if doc_type == 'bill':
        lines += [f"Invoice #: {data.get('invoice_number','INV-001')}", f"Date: {data.get('date', dt.date.today().isoformat())}", f"Bill To: {data.get('client_name','')}", '', f"{'Item':<30}{'Qty':>6}{'Price':>10}{'Total':>10}", '-' * 56]
        total = 0.0
        for it in (data.get('items') or []):
            qty = float(it.get('qty',1)); price = float(it.get('price',0)); line = qty*price; total += line
            lines.append(f"{str(it.get('name',''))[:28]:<30}{qty:>6g}{price:>10.2f}{line:>10.2f}")
        lines += ['-'*56, f"{'Total':>46}{data.get('currency','$')}{total:.2f}"]
    elif doc_type == 'resume':
        if data.get('full_name'): lines.append(data['full_name'])
        if data.get('headline'): lines.append(data['headline'])
        contact = ' • '.join(filter(None, [data.get('email'), data.get('phone'), data.get('location'), data.get('website')]))
        if contact: lines.append(contact)
        for k, label in [('summary','PROFILE'), ('experience','WORK EXPERIENCE'), ('education','EDUCATION'), ('skills','TECHNICAL SKILLS'), ('projects','PROJECTS'), ('hobbies','HOBBIES'), ('languages','LANGUAGES KNOWN')]:
            if data.get(k): lines += ['', label, '-'*len(label), str(data[k])]
    else:
        for sec in (data.get('sections') or []): lines += ['', sec.get('heading',''), '-'*len(sec.get('heading','')), sec.get('body','')]
        if data.get('content'): lines += ['', _expand_with_llm(doc_type, data)]
    with open(out, 'w', encoding='utf-8') as f: f.write('\n'.join(lines))

def _build_md(doc_type, data, out):
    title = data.get('title') or DOC_TEMPLATES.get(doc_type, {}).get('label', 'Document')
    lines = [f"# {title}", '']
    if doc_type == 'bill':
        lines += [f"**Invoice #:** {data.get('invoice_number','INV-001')}  ", f"**Date:** {data.get('date', dt.date.today().isoformat())}  ", f"**Bill To:** {data.get('client_name','')}", '', '| Item | Qty | Price | Total |', '|---|---:|---:|---:|']
        total = 0.0
        for it in (data.get('items') or []):
            qty = float(it.get('qty',1)); price = float(it.get('price',0)); line = qty*price; total += line
            lines.append(f"| {it.get('name','')} | {qty:g} | {price:.2f} | {line:.2f} |")
        lines += ['', f"**Total: {data.get('currency','$')}{total:.2f}**"]
    elif doc_type == 'resume':
        if data.get('full_name'): lines.append(f"## {data['full_name']}")
        if data.get('headline'): lines.append(f"*{data['headline']}*")
        contact = ' • '.join(filter(None, [data.get('email'), data.get('phone'), data.get('location'), data.get('website')]))
        if contact: lines.append(contact)
        for k, label in [('summary','Profile'), ('experience','Work Experience'), ('education','Education'), ('skills','Technical Skills'), ('projects','Projects'), ('hobbies','Hobbies'), ('languages','Languages Known')]:
            if data.get(k): lines += ['', f'### {label}', str(data[k])]
    else:
        for sec in (data.get('sections') or []): lines += ['', f"## {sec.get('heading','')}", sec.get('body','')]
        if data.get('content'): lines += ['', _expand_with_llm(doc_type, data)]
    with open(out, 'w', encoding='utf-8') as f: f.write('\n'.join(lines))

# Blueprint Binding Dictionary Map
BUILDERS = {
    'docx': (_build_docx, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'),
    'pdf':  (_build_pdf,  'application/pdf'),
    'xlsx': (_build_xlsx, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'),
    'pptx': (_build_pptx, 'application/vnd.openxmlformats-officedocument.presentationml.presentation'),
    'txt':  (_build_txt,  'text/plain'),
    'md':   (_build_md,   'text/markdown'),
}